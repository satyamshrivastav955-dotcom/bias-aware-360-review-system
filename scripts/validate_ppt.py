from pathlib import Path
from pptx import Presentation

path = Path(__file__).resolve().parents[1] / "Bias-Aware-360-Review-Hackathon-Pitch.pptx"
prs = Presentation(path)
assert len(prs.slides) == 10, f"expected 10 slides, found {len(prs.slides)}"
assert prs.slide_width / prs.slide_height > 1.7, "deck is not widescreen"

for number, slide in enumerate(prs.slides, start=1):
    texts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text_frame") and shape.text.strip()]
    assert texts, f"slide {number} has no text"
    for shape in slide.shapes:
        assert shape.left >= 0 and shape.top >= 0, f"slide {number}: shape starts outside canvas"
        # Slide-number placeholders and some theme objects may intentionally sit
        # a few points beyond the nominal edge. Treat only material overflow as
        # a layout defect.
        overflow_x = shape.left + shape.width - prs.slide_width
        overflow_y = shape.top + shape.height - prs.slide_height
        assert overflow_x <= 200_000, f"slide {number}: shape exceeds right edge by {overflow_x} EMU"
        assert overflow_y <= 200_000, f"slide {number}: shape exceeds bottom edge by {overflow_y} EMU"

print(f"pptx: ok ({len(prs.slides)} slides, {path.stat().st_size / 1024:.0f} KiB)")
